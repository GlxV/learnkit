from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from app.core.extractors.ocr_extractor import OcrExtractor, SupportsOcr


class PptxExtractor:
    def __init__(self, ocr_extractor: SupportsOcr | None = None) -> None:
        self.ocr_extractor = ocr_extractor or OcrExtractor()

    def extract_text(self, file_path: str | Path) -> str:
        return self.extract(file_path)[0]

    def extract(self, file_path: str | Path) -> tuple[str, list[str], int]:
        presentation = Presentation(str(file_path))
        parts: list[str] = []
        warnings: list[str] = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            slide_parts: list[str] = []
            seen_texts: set[str] = set()
            image_count = 0
            image_text_count = 0
            for shape in self._iter_shapes(slide.shapes):
                if getattr(shape, "has_table", False):
                    rows: list[str] = []
                    for row in shape.table.rows:
                        rows.append(" | ".join(cell.text.strip() for cell in row.cells))
                    if rows:
                        self._append_unique(slide_parts, seen_texts, "Tabela:\n" + "\n".join(rows))
                if hasattr(shape, "text"):
                    text = str(shape.text).strip()
                    if text:
                        self._append_unique(slide_parts, seen_texts, text)
                image_text = self._extract_image_text(shape)
                if image_text is not None:
                    image_count += 1
                    if image_text:
                        image_text_count += 1
                        self._append_unique(slide_parts, seen_texts, "Texto em imagem:\n" + image_text)
            if slide_parts:
                parts.append(f"## Slide {slide_index}\n\n" + "\n\n".join(slide_parts))
            if image_count and image_text_count == 0 and not self.ocr_extractor.available:
                warnings.append(
                    f"Slide {slide_index}: contem imagem, mas OCR local nao esta disponivel."
                )
        return "\n\n".join(parts).strip(), warnings, len(presentation.slides)

    def _iter_shapes(self, shapes) -> list[object]:
        items: list[object] = []
        for shape in shapes:
            items.append(shape)
            children = getattr(shape, "shapes", None)
            if children is not None:
                items.extend(self._iter_shapes(children))
        return items

    def _extract_image_text(self, shape: object) -> str | None:
        try:
            image = shape.image
        except (AttributeError, ValueError):
            return None
        suffix = f".{getattr(image, 'ext', 'png')}"
        try:
            return self.ocr_extractor.extract_image_bytes(image.blob, suffix).strip()
        except Exception:
            return ""

    def _append_unique(self, parts: list[str], seen: set[str], text: str) -> None:
        normalized = " ".join(text.casefold().split())
        if not normalized or normalized in seen:
            return
        seen.add(normalized)
        parts.append(text)
