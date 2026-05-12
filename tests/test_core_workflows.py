from pathlib import Path
from zipfile import ZipFile

import fitz
from pptx import Presentation

from app.core.extractors.file_extractor import FileExtractor
from app.core.importer.ai_response_parser import AIResponseParser
from app.core.prompt.prompt_builder import PromptBuilder
from app.core.services.backup_service import BackupService
from app.core.services.block_service import BlockService
from app.core.services.module_service import ModuleService
from app.core.services.subject_service import SubjectService
from app.core.storage.local_storage import LocalStorage


AI_RESPONSE = """# RESUMO

- Funcoes relacionam entradas e saidas.

# FLASHCARDS

## Card 1

Pergunta: O que e uma funcao?
Resposta: Uma relacao entre entrada e saida.

# PERGUNTAS

## Pergunta 1

Enunciado: Qual alternativa descreve uma funcao?

A) Uma relacao entre entrada e saida
B) Um arquivo
C) Uma tabela de backups
D) Uma imagem

Gabarito: A
Explicacao: A funcao associa valores de entrada a saidas.
"""


def test_end_to_end_block_flow_with_real_pdf_and_pptx(tmp_path: Path) -> None:
    pdf_path = tmp_path / "funcoes.pdf"
    document = fitz.open()
    page = document.new_page()
    page.insert_text((72, 72), "Funcao do primeiro grau tem forma f(x) = ax + b.")
    document.save(pdf_path)
    document.close()

    pptx_path = tmp_path / "slides.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Grafico de uma funcao linear"
    presentation.save(pptx_path)

    storage = LocalStorage(tmp_path / "data")
    SubjectService(storage).create_subject("Matematica")
    ModuleService(storage).create_module("Matematica", "1 Trimestre")
    block_service = BlockService(storage)

    block = block_service.create_block(
        "Matematica",
        "1 Trimestre",
        "Funcoes do 1 grau",
        [pdf_path, pptx_path],
    )
    prompt = block_service.generate_prompt("Matematica", "1 Trimestre", block.title)
    parsed_block = block_service.import_ai_response(
        "Matematica",
        "1 Trimestre",
        block.title,
        AI_RESPONSE,
    )

    assert "Funcao do primeiro grau" in block.extracted_content.text
    assert "Grafico de uma funcao linear" in block.extracted_content.text
    assert "Conteudo fornecido" in prompt
    assert parsed_block.summary is not None
    assert len(parsed_block.flashcards) == 1
    assert len(parsed_block.questions) == 1
    assert (
        tmp_path
        / "data"
        / "subjects"
        / "matematica"
        / "modules"
        / "1_trimestre"
        / "blocks"
        / "funcoes_do_1_grau"
        / "summary.md"
    ).exists()


def test_backup_service_exports_subject_as_zip(tmp_path: Path) -> None:
    storage = LocalStorage(tmp_path / "data")
    SubjectService(storage).create_subject("Historia")
    ModuleService(storage).create_module("Historia", "Revisao Final")
    BlockService(storage).create_block("Historia", "Revisao Final", "Brasil Republica")

    archive_path = BackupService(storage).export_subject("Historia", tmp_path / "backups")

    assert archive_path.exists()
    assert archive_path.suffix == ".zip"
    with ZipFile(archive_path) as archive:
        names = archive.namelist()
    assert "historia/subject.json" in names
    assert "historia/modules/revisao_final/blocks/brasil_republica/block.json" in names


def test_save_imported_package_after_separate_extraction_and_parse(tmp_path: Path) -> None:
    material = tmp_path / "material.md"
    material.write_text("# Funcoes\n\nFuncao linear tem forma f(x) = ax + b.", encoding="utf-8")

    storage = LocalStorage(tmp_path / "data")
    SubjectService(storage).create_subject("Matematica")
    ModuleService(storage).create_module("Matematica", "Prova 1")

    extraction = FileExtractor().extract_files([material])
    prompt = PromptBuilder().build(
        "Matematica",
        "Prova 1",
        "Funcoes lineares",
        extraction.combined_content,
    )
    parsed = AIResponseParser().parse(AI_RESPONSE)

    block = BlockService(storage).save_imported_package(
        "Matematica",
        "Prova 1",
        "Funcoes lineares",
        extraction,
        prompt,
        AI_RESPONSE,
        parsed,
    )

    _, _, loaded = storage.get_block_by_id(block.id)
    assert "Funcao linear" in loaded.extracted_content.text
    assert "Conteudo fornecido" in loaded.generated_prompt
    assert loaded.ai_response_raw == AI_RESPONSE
    assert loaded.summary is not None
    assert len(loaded.flashcards) == 1
    assert len(loaded.questions) == 1
    assert (
        tmp_path
        / "data"
        / "subjects"
        / "matematica"
        / "modules"
        / "prova_1"
        / "blocks"
        / "funcoes_lineares"
        / "progress.json"
    ).exists()
