from pathlib import Path

import fitz
from PIL import Image, ImageDraw
from pptx import Presentation

from app.core.extractors.file_extractor import FileExtractor
from app.core.extractors.pptx_extractor import PptxExtractor


def test_extract_pdf_and_pptx_without_stopping_on_unsupported_file(tmp_path: Path) -> None:
    pdf_path = tmp_path / "aula.pdf"
    pdf_doc = fitz.open()
    page = pdf_doc.new_page()
    page.insert_text((72, 72), "Conteudo do PDF")
    pdf_doc.save(pdf_path)
    pdf_doc.close()

    pptx_path = tmp_path / "slides.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Conteudo do PPTX"
    presentation.save(pptx_path)

    unsupported_path = tmp_path / "notas.xlsx"
    unsupported_path.write_text("nao suportado", encoding="utf-8")

    result = FileExtractor().extract_files([pdf_path, pptx_path, unsupported_path])

    successful = [item for item in result.files if item.error_message is None]
    failed = [item for item in result.files if item.error_message is not None]

    assert len(successful) == 2
    assert "Conteudo do PDF" in result.combined_content.text
    assert "Conteudo do PPTX" in result.combined_content.text
    assert len(failed) == 1
    assert "nao suportado" in failed[0].error_message.lower()


def test_export_extracted_text_as_markdown_and_txt(tmp_path: Path) -> None:
    extractor = FileExtractor()
    extraction = extractor.extract_files([])
    extraction.file_texts = {"Aula_01.pdf": "Texto 1", "Aula_02.pptx": "Texto 2"}

    md_path = extractor.export_extracted_text(extraction, tmp_path / "saida.md", "md")
    txt_path = extractor.export_extracted_text(extraction, tmp_path / "saida.txt", "txt")

    assert md_path.read_text(encoding="utf-8").startswith("# Conteudo extraido")
    assert "## Arquivo: Aula_01.pdf" in md_path.read_text(encoding="utf-8")
    assert "Arquivo: Aula_02.pptx" in txt_path.read_text(encoding="utf-8")


def test_file_extractor_reads_javascript_as_plain_text(tmp_path: Path) -> None:
    js_path = tmp_path / "pilha.js"
    js_path.write_text("class Stack { push(item) { return item } }", encoding="utf-8")

    result = FileExtractor().extract_files([js_path])

    assert result.files[0].error_message is None
    assert "class Stack" in result.combined_content.text


def test_pptx_extractor_sends_embedded_images_to_ocr(tmp_path: Path) -> None:
    image_path = tmp_path / "slide_image.png"
    image = Image.new("RGB", (480, 160), "white")
    draw = ImageDraw.Draw(image)
    draw.text((24, 56), "Texto dentro da imagem", fill="black")
    image.save(image_path)

    pptx_path = tmp_path / "imagem.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(str(image_path), 0, 0)
    presentation.save(pptx_path)

    class FakeOcr:
        available = True

        def extract_image_bytes(self, image_bytes: bytes, suffix: str = ".png") -> str:
            assert image_bytes
            assert suffix == ".png"
            return "Texto OCR da imagem"

    text, warnings, slide_count = PptxExtractor(ocr_extractor=FakeOcr()).extract(pptx_path)

    assert slide_count == 1
    assert "Texto em imagem" in text
    assert "Texto OCR da imagem" in text
    assert warnings == []
